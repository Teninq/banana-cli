#!/usr/bin/env python
"""
export_pptx_from_descriptions.py - Convert API descriptions to a styled PPTX
"""
import sys
import json
import re
import argparse
from pathlib import Path

_ROOT = Path(__file__).parent
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))

import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ─── Color Palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT    = RGBColor(0x1E, 0x88, 0xE5)   # blue
TITLE_COL = RGBColor(0xFF, 0xFF, 0xFF)   # white
BODY_COL  = RGBColor(0xE0, 0xE8, 0xF0)  # light grey-blue
PART_COL  = RGBColor(0x90, 0xCA, 0xF9)  # light blue (for section label)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_tf_text(tf, text, size_pt, bold=False, color=BODY_COL, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes._spTree.clear()

    # Background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)
    # Accent bar left
    add_rect(slide, 0, 0, Inches(0.15), SLIDE_H, ACCENT)
    # Accent line
    add_rect(slide, Inches(0.15), Inches(3.1), Inches(8), Emu(4), ACCENT)

    # Title text box
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(11.5), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = TITLE_COL

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(11.5), Inches(1.0))
        set_tf_text(tb2.text_frame, subtitle, 18, color=PART_COL)

    return slide


def add_content_slide(prs, title, bullets, part=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes._spTree.clear()

    # Background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)
    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), ACCENT)
    # Left accent bar
    add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, ACCENT)

    y_start = Inches(0.22)

    # Section label (part)
    if part:
        tb_p = slide.shapes.add_textbox(Inches(0.3), y_start, Inches(12.5), Inches(0.4))
        set_tf_text(tb_p.text_frame, part, 11, color=PART_COL, align=PP_ALIGN.LEFT)
        y_start += Inches(0.42)

    # Title
    tb_t = slide.shapes.add_textbox(Inches(0.3), y_start, Inches(12.5), Inches(1.0))
    set_tf_text(tb_t.text_frame, title, 28, bold=True, color=TITLE_COL)
    y_start += Inches(1.1)

    # Accent underline
    add_rect(slide, Inches(0.3), y_start, Inches(1.5), Emu(3), ACCENT)
    y_start += Inches(0.18)

    # Bullet points
    if bullets:
        tb_b = slide.shapes.add_textbox(Inches(0.45), y_start, Inches(12.4), Inches(7.5) - y_start - Inches(0.2))
        tf = tb_b.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if not bullet.strip():
                continue
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4)
            p.space_after = Pt(2)
            run = p.add_run()
            # strip leading bullet chars
            text = bullet.lstrip("•·-* ").strip()
            run.text = f"  {text}"
            run.font.size = Pt(15)
            run.font.color.rgb = BODY_COL

    return slide


def extract_bullets_from_description(desc_content):
    """Extract bullet points from description content dict or string."""
    if not desc_content:
        return []

    text = ""
    if isinstance(desc_content, dict):
        text = desc_content.get("text", "")
    elif isinstance(desc_content, str):
        try:
            parsed = json.loads(desc_content)
            text = parsed.get("text", desc_content)
        except Exception:
            text = desc_content

    if not text:
        return []

    lines = text.splitlines()
    bullets = []
    in_body = False
    for line in lines:
        raw = line.strip()
        if not raw:
            continue
        # Skip section headers
        if raw.startswith("#") or raw.startswith("页面类型") or raw.startswith("图片需求"):
            in_body = False
            continue
        # Detect body start after header
        if raw.startswith("*") or re.match(r"^\d+[\.\)、]", raw) or raw.startswith("•") or raw.startswith("-"):
            in_body = True

        if in_body and raw not in ("*", "-", "•"):
            # Clean markdown formatting
            clean = re.sub(r"\*\*(.+?)\*\*", r"\1", raw)  # bold
            clean = re.sub(r"\*(.+?)\*", r"\1", clean)  # italic
            clean = clean.lstrip("* •·-").strip()
            # Skip too-short lines
            if len(clean) < 4:
                continue
            bullets.append(clean)

    return bullets[:7]  # max 7 bullet points per slide


def fetch_project_pages(base_url, project_id):
    resp = requests.get(f"{base_url}/api/projects/{project_id}", timeout=30)
    resp.raise_for_status()
    data = resp.json().get("data", {})
    return data.get("pages", [])


def build_pptx(pages, output_path, title="Presentation"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title slide
    add_title_slide(prs, title, subtitle="")

    for page in sorted(pages, key=lambda p: p.get("order_index", 0)):
        oc = page.get("outline_content") or {}
        if isinstance(oc, str):
            try:
                oc = json.loads(oc)
            except Exception:
                oc = {}

        slide_title = oc.get("title", "")
        if not slide_title:
            continue

        # Try to get bullets: first from outline_content.points, then from description
        bullets = oc.get("points", []) or []
        if not bullets:
            bullets = extract_bullets_from_description(page.get("description_content"))

        part = page.get("part", "")
        if isinstance(part, str) and part.strip():
            part_text = part.strip()
        else:
            part_text = ""

        add_content_slide(prs, slide_title, bullets, part=part_text)

    prs.save(str(output_path))
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Export project descriptions to PPTX")
    parser.add_argument("project_id", help="Project ID")
    parser.add_argument("--url", default="http://localhost:5000", help="Backend URL")
    parser.add_argument("--title", default="", help="Presentation title")
    parser.add_argument("--out", default="", help="Output file path (.pptx)")
    args = parser.parse_args()

    print(f"Fetching project {args.project_id[:8]}...")
    pages = fetch_project_pages(args.url, args.project_id)
    print(f"Found {len(pages)} pages")

    out_path = Path(args.out) if args.out else Path(f"{args.project_id[:8]}.pptx")
    title = args.title or (pages and (pages[0].get("outline_content") or {}).get("title", "")) or "Presentation"

    build_pptx(pages, out_path, title=title)
    print(f"[DONE] Saved: {out_path.absolute()}")


if __name__ == "__main__":
    main()
