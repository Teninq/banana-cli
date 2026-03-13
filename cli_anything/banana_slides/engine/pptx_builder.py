"""
PPTXBuilder — simplified builder for editable PPTX slides.

Ported from backend/utils/pptx_builder.py with reduced scope:
  - No HTML table parsing
  - No custom font file dependency (uses estimation for font sizing)
  - Focused on add_text_element / add_image_element with bbox positioning
"""
import io
import logging
import os
from datetime import datetime, timezone
from typing import Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

# ── Constants ────────────────────────────────────────────────────────────
DEFAULT_SLIDE_WIDTH = 10.0   # inches (16:9)
DEFAULT_SLIDE_HEIGHT = 5.625
DEFAULT_DPI = 96
MIN_FONT_SIZE = 6    # pt
MAX_FONT_SIZE = 200  # pt

ALIGNMENT_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


class PPTXBuilder:
    """Builds editable PPTX presentations with precise element positioning."""

    def __init__(self, slide_width_inches: float = DEFAULT_SLIDE_WIDTH,
                 slide_height_inches: float = DEFAULT_SLIDE_HEIGHT):
        self.slide_width = slide_width_inches
        self.slide_height = slide_height_inches
        self.prs: Optional[Presentation] = None

    def create_presentation(self) -> Presentation:
        """Create a new blank presentation with configured dimensions."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_width)
        self.prs.slide_height = Inches(self.slide_height)
        try:
            core = self.prs.core_properties
            now = datetime.now(timezone.utc)
            core.author = "banana-slides"
            core.last_modified_by = "banana-slides"
            core.created = now
            core.modified = now
        except Exception:
            pass
        return self.prs

    def setup_presentation_size(self, width_pixels: int, height_pixels: int,
                                dpi: int = DEFAULT_DPI) -> None:
        """Set slide dimensions from pixel sizes."""
        w_inches = width_pixels / dpi
        h_inches = height_pixels / dpi
        # Clamp to python-pptx limits (1-56 inches) preserving aspect ratio
        max_dim = max(w_inches, h_inches)
        if max_dim > 56:
            scale = 56 / max_dim
            w_inches *= scale
            h_inches *= scale
            logger.warning("Scaled down to %.2f x %.2f inches", w_inches, h_inches)
        self.slide_width = max(1.0, w_inches)
        self.slide_height = max(1.0, h_inches)
        if self.prs:
            self.prs.slide_width = Inches(self.slide_width)
            self.prs.slide_height = Inches(self.slide_height)

    def add_blank_slide(self):
        """Add a blank slide and return it."""
        if not self.prs:
            self.create_presentation()
        blank_layout = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(blank_layout)

    @staticmethod
    def pixels_to_inches(pixels: float, dpi: int = DEFAULT_DPI) -> float:
        return pixels / dpi

    def add_text_element(self, slide, text: str, bbox: Tuple[float, float, float, float],
                         dpi: int = DEFAULT_DPI, align: str = "left",
                         text_style: Optional[dict] = None) -> None:
        """
        Add a text box at the given bbox position.

        Args:
            slide: pptx slide object
            text: The text content
            bbox: (x0, y0, x1, y1) in pixels
            dpi: Resolution for pixel-to-inch conversion
            align: Text alignment (left/center/right/justify)
            text_style: Optional dict with keys: color, bold, italic, font_size, font_name
        """
        style = text_style or {}
        x0, y0, x1, y1 = bbox
        # Expand by 1% to prevent text overflow
        w_px = x1 - x0
        h_px = y1 - y0
        x0 -= w_px * 0.005
        y0 -= h_px * 0.005
        w_px *= 1.01
        h_px *= 1.01

        left = Inches(x0 / dpi)
        top = Inches(y0 / dpi)
        width = Inches(w_px / dpi)
        height = Inches(h_px / dpi)

        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        # Remove internal margins for tight positioning
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

        # Font size: use explicit value, or calculate from bbox
        font_size = style.get("font_size")
        if not font_size:
            font_size = self.calculate_font_size(bbox, text, dpi=dpi)

        p = tf.paragraphs[0]
        p.alignment = ALIGNMENT_MAP.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text.replace("·", "•")
        run.font.size = Pt(font_size)
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)

        if "font_name" in style:
            run.font.name = style["font_name"]

        color = style.get("color")
        if color:
            if isinstance(color, str) and color.startswith("#") and len(color) == 7:
                r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
                run.font.color.rgb = RGBColor(r, g, b)
            elif isinstance(color, (list, tuple)) and len(color) == 3:
                run.font.color.rgb = RGBColor(*color)

    def add_image_element(self, slide, image_path: str,
                          bbox: Tuple[float, float, float, float],
                          dpi: int = DEFAULT_DPI) -> None:
        """Add an image at the given bbox position."""
        if not os.path.exists(image_path):
            logger.warning("Image not found for element: %s", image_path)
            return
        x0, y0, x1, y1 = bbox
        left = Inches(x0 / dpi)
        top = Inches(y0 / dpi)
        width = Inches((x1 - x0) / dpi)
        height = Inches((y1 - y0) / dpi)
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            logger.error("Failed to add image element: %s", e)

    def calculate_font_size(self, bbox: Tuple[float, float, float, float],
                            text: str, text_level: int = 0,
                            dpi: int = DEFAULT_DPI) -> int:
        """
        Estimate the best font size for text within a bounding box.

        Uses a simple estimation based on character count, bbox dimensions,
        and CJK character width ratios (no font file dependency).
        """
        x0, y0, x1, y1 = bbox
        w_inches = (x1 - x0) / dpi
        h_inches = (y1 - y0) / dpi
        w_pt = w_inches * 72
        h_pt = h_inches * 72

        if not text.strip():
            return 12

        # Estimate char widths: CJK chars ≈ 1.0x font_size, others ≈ 0.55x
        def _effective_chars(t: str) -> float:
            count = 0.0
            for ch in t:
                if ord(ch) > 0x2E80:
                    count += 1.0
                else:
                    count += 0.55
            return max(count, 1.0)

        lines = text.split("\n")
        num_lines = len(lines)
        max_line_chars = max(_effective_chars(line) for line in lines)

        # Binary search for largest font that fits
        lo, hi = MIN_FONT_SIZE, MAX_FONT_SIZE
        best = MIN_FONT_SIZE
        while lo <= hi:
            mid = (lo + hi) // 2
            line_height = mid * 1.2
            chars_per_line = w_pt / (mid * 0.6)  # approximate
            # Estimate total lines needed (with wrapping)
            total_lines = 0
            for line in lines:
                ec = _effective_chars(line)
                needed = max(1, int((ec * mid) / w_pt + 0.99))
                total_lines += needed
            total_height = total_lines * line_height
            if total_height <= h_pt and max_line_chars * mid <= w_pt * 1.1:
                best = mid
                lo = mid + 1
            else:
                hi = mid - 1

        return max(MIN_FONT_SIZE, min(best, MAX_FONT_SIZE))

    def save(self, output_path: str) -> str:
        """Save the presentation to disk."""
        if not self.prs:
            raise RuntimeError("No presentation to save. Call create_presentation() first.")
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        self.prs.save(output_path)
        logger.info("Saved PPTX to %s", output_path)
        return output_path

    def to_bytes(self) -> bytes:
        """Return the presentation as bytes."""
        if not self.prs:
            raise RuntimeError("No presentation to export.")
        buf = io.BytesIO()
        self.prs.save(buf)
        return buf.getvalue()

    def get_presentation(self) -> Optional[Presentation]:
        return self.prs
