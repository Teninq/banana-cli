"""
PPTX export utilities (standalone, no Flask)

Two export modes:
  1. Text PPTX -- styled slides with text content (dark theme)
  2. Image PPTX -- full-bleed slide images
"""
import math
import logging
import os
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# ── Color palette (dark theme) ────────────────────────────────────────────
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT = RGBColor(0x1E, 0x88, 0xE5)
TITLE_COL = RGBColor(0xFF, 0xFF, 0xFF)
BODY_COL = RGBColor(0xE0, 0xE8, 0xF0)
PART_COL = RGBColor(0x90, 0xCA, 0xF9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _get_page_size_inches(aspect_ratio: str = '16:9',
                          base: float = 10.0) -> Tuple[float, float]:
    try:
        w, h = (float(x) for x in aspect_ratio.split(':'))
        if not (math.isfinite(w) and math.isfinite(h) and w > 0 and h > 0):
            raise ValueError
    except (ValueError, AttributeError):
        w, h = 16.0, 9.0
    if w >= h:
        return base, base * h / w
    return base * w / h, base


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _set_tf_text(tf, text, size_pt, bold=False, color=BODY_COL, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


# ── Text PPTX ─────────────────────────────────────────────────────────────

def export_text_pptx(pages: List[dict], output_path: str,
                     title: str = "Presentation") -> str:
    """Build a styled text-only PPTX from page data dicts.

    Each page dict should have:
        outline_content: {title, points}
        part: str (optional)
        description_content: str or {text: str} (optional, for bullet extraction)
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title slide
    _add_title_slide(prs, title)

    for page in sorted(pages, key=lambda p: p.get("order_index", 0)):
        oc = page.get("outline_content") or {}
        slide_title = oc.get("title", "")
        if not slide_title:
            continue
        bullets = oc.get("points", []) or []
        if not bullets:
            bullets = _extract_bullets(page.get("description_content"))
        part_text = (page.get("part") or "").strip()
        _add_content_slide(prs, slide_title, bullets, part=part_text)

    prs.save(str(output_path))
    return str(output_path)


def _add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes._spTree.clear()
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)
    _add_rect(slide, 0, 0, Inches(0.15), SLIDE_H, ACCENT)
    _add_rect(slide, Inches(0.15), Inches(3.1), Inches(8), Emu(4), ACCENT)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(11.5), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = TITLE_COL

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(11.5), Inches(1.0))
        _set_tf_text(tb2.text_frame, subtitle, 18, color=PART_COL)


def _add_content_slide(prs, title, bullets, part=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes._spTree.clear()
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), ACCENT)
    _add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, ACCENT)

    y = Inches(0.22)
    if part:
        tb_p = slide.shapes.add_textbox(Inches(0.3), y, Inches(12.5), Inches(0.4))
        _set_tf_text(tb_p.text_frame, part, 11, color=PART_COL)
        y += Inches(0.42)

    tb_t = slide.shapes.add_textbox(Inches(0.3), y, Inches(12.5), Inches(1.0))
    _set_tf_text(tb_t.text_frame, title, 28, bold=True, color=TITLE_COL)
    y += Inches(1.1)
    _add_rect(slide, Inches(0.3), y, Inches(1.5), Emu(3), ACCENT)
    y += Inches(0.18)

    if bullets:
        tb_b = slide.shapes.add_textbox(
            Inches(0.45), y, Inches(12.4), Inches(7.5) - y - Inches(0.2))
        tf = tb_b.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if not bullet.strip():
                continue
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(4)
            p.space_after = Pt(2)
            run = p.add_run()
            text = bullet.lstrip("•·-* ").strip()
            run.text = f"  {text}"
            run.font.size = Pt(15)
            run.font.color.rgb = BODY_COL


def _extract_bullets(desc_content) -> List[str]:
    """Extract bullet points from description content."""
    import re
    if not desc_content:
        return []
    text = ""
    if isinstance(desc_content, dict):
        text = desc_content.get("text", "")
    elif isinstance(desc_content, str):
        text = desc_content
    if not text:
        return []

    bullets = []
    in_body = False
    for line in text.splitlines():
        raw = line.strip()
        if not raw:
            continue
        if raw.startswith("#") or raw.startswith("页面类型"):
            in_body = False
            continue
        if raw.startswith("*") or raw.startswith("•") or raw.startswith("-") or (len(raw) > 1 and raw[0].isdigit()):
            in_body = True
        if in_body and raw not in ("*", "-", "•"):
            clean = re.sub(r"\*\*(.+?)\*\*", r"\1", raw)
            clean = re.sub(r"\*(.+?)\*", r"\1", clean)
            clean = clean.lstrip("* •·-").strip()
            if len(clean) >= 4:
                bullets.append(clean)
    return bullets[:7]


# ── Image PPTX ────────────────────────────────────────────────────────────

def export_image_pptx(image_paths: List[str], output_path: str,
                      aspect_ratio: str = "16:9") -> str:
    """Create PPTX where each slide is a full-bleed image."""
    prs = Presentation()

    try:
        core = prs.core_properties
        now = datetime.now(timezone.utc)
        core.author = "banana-slides"
        core.last_modified_by = "banana-slides"
        core.created = now
        core.modified = now
    except Exception:
        pass

    page_w, page_h = _get_page_size_inches(aspect_ratio)
    prs.slide_width = Inches(page_w)
    prs.slide_height = Inches(page_h)

    for image_path in image_paths:
        if not os.path.exists(image_path):
            logger.warning("Image not found: %s", image_path)
            continue
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            image_path, left=0, top=0,
            width=prs.slide_width, height=prs.slide_height,
        )

    prs.save(str(output_path))
    return str(output_path)
